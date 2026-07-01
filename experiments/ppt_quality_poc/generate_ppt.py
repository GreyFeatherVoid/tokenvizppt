#!/usr/bin/env python3
"""Standalone PPT quality proof-of-concept.

This script intentionally does not import the tokenvizPPT backend. It is a small,
direct pipeline for testing whether document understanding + AI planning +
optional image generation + PPTX rendering can produce a stronger deck.
"""

from __future__ import annotations

import argparse
import base64
import csv
import json
import os
import re
import sys
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Any


ROOT = Path(__file__).resolve().parents[2]
DEFAULT_ENV = ROOT / "backend" / ".env"


@dataclass(frozen=True)
class ModelConfig:
    model: str
    api_key: str
    base_url: str | None
    temperature: float


@dataclass(frozen=True)
class ImageConfig:
    enabled: bool
    model: str
    api_key: str
    base_url: str | None
    size: str


def main() -> None:
    args = parse_args()
    from dotenv import load_dotenv
    from openai import OpenAI

    load_dotenv(DEFAULT_ENV)
    load_dotenv(Path(".env"), override=True)

    input_path = Path(args.input).expanduser().resolve()
    output_path = Path(args.output).expanduser().resolve()
    template_path = Path(args.template).expanduser().resolve() if args.template else None
    output_path.parent.mkdir(parents=True, exist_ok=True)

    text = parse_document(input_path)
    if not text.strip():
        raise SystemExit(f"No readable text found in {input_path}")
    text = compact_text(text, args.max_chars)

    llm_config = load_llm_config(args)
    client = OpenAI(api_key=llm_config.api_key, base_url=llm_config.base_url)

    print(f"[poc] parsed {len(text):,} chars from {input_path.name}")
    print("[poc] asking model for document intelligence...")
    intelligence = ask_json(
        client,
        llm_config,
        system=(
            "You are a principal strategy analyst. Read source material deeply and "
            "return precise, structured JSON. No markdown."
        ),
        user=build_intelligence_prompt(args.topic, text),
    )

    print("[poc] asking model for deck plan...")
    deck_plan = ask_json(
        client,
        llm_config,
        system=(
            "You are an elite presentation strategist and art director. Return only "
            "valid JSON for a practical, visually rich PowerPoint deck. No markdown."
        ),
        user=build_deck_prompt(
            topic=args.topic,
            slide_count=args.slides,
            intelligence=intelligence,
            template_note=template_note(template_path),
        ),
    )
    deck_plan = normalize_deck_plan(deck_plan, args.topic, args.slides)

    image_config = load_image_config(args)
    image_paths: dict[int, Path] = {}
    if args.with_images and image_config.enabled:
        print("[poc] generating AI images...")
        image_paths = generate_images(deck_plan, image_config, output_path.parent)
    elif args.with_images:
        print("[poc] image generation requested, but image config is not enabled")

    print("[poc] rendering PPTX...")
    render_pptx(deck_plan, output_path, template_path, image_paths)
    print(f"[poc] wrote {output_path}")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a standalone PPTX quality POC.")
    parser.add_argument("--input", required=True, help="Input txt/md/csv/pdf/docx file.")
    parser.add_argument("--output", default="outputs/poc_deck.pptx", help="Output PPTX path.")
    parser.add_argument("--template", help="Optional PPTX template path.")
    parser.add_argument("--topic", default="Document briefing", help="Deck topic.")
    parser.add_argument("--slides", type=int, default=8, help="Number of slides, 3-20.")
    parser.add_argument("--max-chars", type=int, default=50000, help="Max parsed chars sent to LLM.")
    parser.add_argument("--with-images", action="store_true", help="Generate and insert AI images.")
    parser.add_argument("--llm-model", help="Override TOKENVIZPPT_LLM_MODEL.")
    parser.add_argument("--llm-base-url", help="Override TOKENVIZPPT_LLM_BASE_URL.")
    parser.add_argument("--image-model", help="Override TOKENVIZPPT_AI_IMAGE_MODEL.")
    parser.add_argument("--image-base-url", help="Override TOKENVIZPPT_AI_IMAGE_BASE_URL.")
    parser.add_argument("--image-size", default="1536x1024", help="AI image size.")
    args = parser.parse_args()
    args.slides = max(3, min(20, int(args.slides)))
    return args


def load_llm_config(args: argparse.Namespace) -> ModelConfig:
    model = args.llm_model or os.getenv("TOKENVIZPPT_LLM_MODEL", "")
    api_key = os.getenv("TOKENVIZPPT_LLM_API_KEY", "")
    base_url = args.llm_base_url or os.getenv("TOKENVIZPPT_LLM_BASE_URL", "") or None
    temperature = float(os.getenv("TOKENVIZPPT_LLM_TEMPERATURE", "0.55"))
    if not model or not api_key:
        raise SystemExit("Missing TOKENVIZPPT_LLM_MODEL or TOKENVIZPPT_LLM_API_KEY")
    return ModelConfig(model=model, api_key=api_key, base_url=base_url, temperature=temperature)


def load_image_config(args: argparse.Namespace) -> ImageConfig:
    enabled = os.getenv("TOKENVIZPPT_AI_IMAGE_ENABLED", "false").lower() == "true"
    model = args.image_model or os.getenv("TOKENVIZPPT_AI_IMAGE_MODEL", "")
    api_key = os.getenv("TOKENVIZPPT_AI_IMAGE_API_KEY", "") or os.getenv("TOKENVIZPPT_LLM_API_KEY", "")
    base_url = args.image_base_url or os.getenv("TOKENVIZPPT_AI_IMAGE_BASE_URL", "") or None
    return ImageConfig(
        enabled=enabled and bool(model and api_key),
        model=model,
        api_key=api_key,
        base_url=base_url,
        size=args.image_size,
    )


def parse_document(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".markdown"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".csv":
        return parse_csv(path)
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix == ".docx":
        return parse_docx(path)
    raise SystemExit(f"Unsupported input type: {suffix}")


def parse_csv(path: Path) -> str:
    lines = []
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for index, row in enumerate(reader):
            lines.append(" | ".join(cell.strip() for cell in row))
            if index >= 500:
                lines.append("[CSV truncated after 500 rows]")
                break
    return "\n".join(lines)


def parse_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages.append(f"\n\n[Page {index + 1}]\n{text}")
    return "\n".join(pages)


def parse_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    chunks: list[str] = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            chunks.append(paragraph.text.strip())
    for table_index, table in enumerate(doc.tables, start=1):
        chunks.append(f"\n[Table {table_index}]")
        for row in table.rows:
            chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(chunks)


def compact_text(text: str, max_chars: int) -> str:
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if len(text) <= max_chars:
        return text
    head = text[: int(max_chars * 0.65)]
    tail = text[-int(max_chars * 0.35) :]
    return f"{head}\n\n[...middle truncated for length...]\n\n{tail}"


def ask_json(client: OpenAI, config: ModelConfig, *, system: str, user: str) -> dict[str, Any]:
    response = client.chat.completions.create(
        model=config.model,
        temperature=config.temperature,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
    )
    content = response.choices[0].message.content or ""
    return json.loads(extract_json(content))


def extract_json(content: str) -> str:
    cleaned = content.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?", "", cleaned).strip()
        cleaned = re.sub(r"```$", "", cleaned).strip()
    start = cleaned.find("{")
    end = cleaned.rfind("}")
    if start == -1 or end == -1 or end <= start:
        raise ValueError(f"Model did not return JSON: {content[:500]}")
    return cleaned[start : end + 1]


def build_intelligence_prompt(topic: str, text: str) -> str:
    return f"""
Topic:
{topic}

Source material:
{text}

Return JSON with:
{{
  "executive_summary": "one paragraph",
  "audience": ["who should hear this deck"],
  "core_message": "the main argument",
  "document_structure": [
    {{"section": "name", "purpose": "why it matters", "key_points": ["..."]}}
  ],
  "key_evidence": [
    {{"claim": "claim", "evidence": "source evidence", "implication": "so what"}}
  ],
  "data_points": [
    {{"label": "metric or fact", "value": "value", "meaning": "interpretation"}}
  ],
  "risks": ["specific risks"],
  "opportunities": ["specific opportunities"],
  "recommended_storyline": ["chapter by chapter narrative"]
}}

Rules:
- Do not copy long paragraphs.
- Prefer concrete claims and evidence over generic summaries.
- If data is missing, say what is missing.
""".strip()


def build_deck_prompt(
    *,
    topic: str,
    slide_count: int,
    intelligence: dict[str, Any],
    template_note: str,
) -> str:
    return f"""
Create a {slide_count}-slide PowerPoint plan for:
{topic}

Document intelligence:
{json.dumps(intelligence, ensure_ascii=False, indent=2)}

Template context:
{template_note}

Return JSON:
{{
  "deck_title": "short title",
  "subtitle": "short subtitle",
  "design_system": {{
    "tone": "visual tone",
    "primary_color": "#hex",
    "accent_color": "#hex",
    "background_color": "#hex",
    "font": "Aptos"
  }},
  "slides": [
    {{
      "layout": "cover|section|insight|comparison|timeline|data|quote|closing",
      "title": "slide title",
      "kicker": "short section label",
      "narrative": "speaker-facing point of the slide",
      "bullets": ["3 to 5 concise bullets"],
      "evidence": "specific evidence or data",
      "visual_type": "image|diagram|cards|chart|none",
      "image_prompt": "prompt for a clean business illustration or empty string",
      "chart": {{
        "type": "bar|line|none",
        "labels": ["A", "B"],
        "values": [1, 2],
        "caption": "what the chart means"
      }}
    }}
  ]
}}

Rules:
- Return exactly {slide_count} slides.
- Make every slide say one clear thing.
- Avoid dense paragraphs.
- Use images only where they clarify the story.
- Use concrete evidence from the source.
- Keep visual prompts specific, non-photorealistic, and presentation-safe.
- {template_note}
""".strip()


def template_note(template_path: Path | None) -> str:
    if not template_path:
        return "No template provided. Use a clean, modern consulting style."
    return (
        f"Template provided: {template_path.name}. Reuse its page size, theme, "
        "master defaults, and existing brand feel where python-pptx allows."
    )


def normalize_deck_plan(plan: dict[str, Any], topic: str, slide_count: int) -> dict[str, Any]:
    slides = plan.get("slides")
    if not isinstance(slides, list):
        slides = []
    while len(slides) < slide_count:
        slides.append(
            {
                "layout": "insight",
                "title": f"{topic} #{len(slides) + 1}",
                "kicker": "Insight",
                "narrative": "A focused supporting point.",
                "bullets": ["Key point", "Evidence", "Implication"],
                "evidence": "",
                "visual_type": "cards",
                "image_prompt": "",
                "chart": {"type": "none", "labels": [], "values": [], "caption": ""},
            }
        )
    plan["slides"] = slides[:slide_count]
    plan.setdefault("deck_title", topic)
    plan.setdefault("subtitle", "AI-generated strategic briefing")
    plan.setdefault(
        "design_system",
        {
            "tone": "clean executive",
            "primary_color": "#183A37",
            "accent_color": "#D98E04",
            "background_color": "#F7F3EA",
            "font": "Aptos",
        },
    )
    return plan


def generate_images(plan: dict[str, Any], config: ImageConfig, output_dir: Path) -> dict[int, Path]:
    image_dir = output_dir / "generated_images"
    image_dir.mkdir(parents=True, exist_ok=True)
    paths: dict[int, Path] = {}
    for index, slide in enumerate(plan["slides"], start=1):
        prompt = str(slide.get("image_prompt") or "").strip()
        if not prompt or str(slide.get("visual_type", "")).lower() not in {"image", "diagram"}:
            continue
        print(f"[poc] image {index}: {prompt[:80]}")
        try:
            image_bytes = request_image(config, prompt)
        except Exception as exc:
            print(f"[poc] image {index} failed: {exc}")
            continue
        path = image_dir / f"slide_{index:02d}.png"
        path.write_bytes(image_bytes)
        paths[index] = path
    return paths


def request_image(config: ImageConfig, prompt: str) -> bytes:
    import requests

    base_url = (config.base_url or "https://api.openai.com/v1").rstrip("/")
    response = requests.post(
        f"{base_url}/images/generations",
        headers={
            "Authorization": f"Bearer {config.api_key}",
            "Content-Type": "application/json",
        },
        json={
            "model": config.model,
            "prompt": prompt,
            "n": 1,
            "size": config.size,
            "response_format": "b64_json",
        },
        timeout=240,
    )
    response.raise_for_status()
    payload = response.json()
    data = payload.get("data") or []
    if not data or not data[0].get("b64_json"):
        raise ValueError("Image response missing b64_json")
    return base64.b64decode(data[0]["b64_json"])


def render_pptx(
    plan: dict[str, Any],
    output_path: Path,
    template_path: Path | None,
    image_paths: dict[int, Path],
) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    prs = Presentation(str(template_path)) if template_path and template_path.exists() else Presentation()
    ensure_wide(prs)
    clear_existing_slides(prs)

    design = plan["design_system"]
    colors = {
        "primary": parse_hex(design.get("primary_color"), RGBColor(24, 58, 55)),
        "accent": parse_hex(design.get("accent_color"), RGBColor(217, 142, 4)),
        "background": parse_hex(design.get("background_color"), RGBColor(247, 243, 234)),
        "ink": RGBColor(30, 35, 38),
        "muted": RGBColor(92, 100, 106),
        "white": RGBColor(255, 255, 255),
    }
    font = str(design.get("font") or "Aptos")

    for index, slide_data in enumerate(plan["slides"], start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        paint_background(slide, prs, colors["background"])
        add_page_number(slide, prs, index, len(plan["slides"]), colors, font)
        layout = str(slide_data.get("layout") or "insight").lower()
        if index == 1 or layout == "cover":
            render_cover(slide, prs, plan, slide_data, colors, font)
        elif layout == "section":
            render_section(slide, prs, slide_data, colors, font)
        elif layout in {"comparison", "timeline"}:
            render_split(slide, prs, slide_data, colors, font, image_paths.get(index))
        elif layout == "data" and has_chart(slide_data):
            render_data_slide(slide, prs, slide_data, colors, font)
        else:
            render_insight(slide, prs, slide_data, colors, font, image_paths.get(index))

    prs.save(output_path)


def ensure_wide(prs: Presentation) -> None:
    from pptx.util import Inches

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def clear_existing_slides(prs: Presentation) -> None:
    # Keep template masters/theme, remove starter slides.
    slide_ids = list(prs.slides._sldIdLst)  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in slide_ids:
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)  # noqa: SLF001


def paint_background(slide, prs: Presentation, color: RGBColor) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def render_cover(slide, prs: Presentation, plan: dict[str, Any], data: dict[str, Any], colors: dict, font: str) -> None:
    add_accent_bar(slide, prs, colors)
    add_text(slide, data.get("kicker") or "Strategic briefing", 0.75, 0.65, 3.8, 0.35, 15, colors["accent"], font, bold=True)
    add_text(slide, plan["deck_title"], 0.72, 1.35, 7.8, 1.8, 42, colors["primary"], font, bold=True)
    add_text(slide, plan.get("subtitle") or data.get("narrative") or "", 0.78, 3.15, 5.9, 0.7, 18, colors["muted"], font)
    add_big_number(slide, prs, "POC", colors, font)
    bullets = data.get("bullets") or []
    add_cards(slide, bullets[:3], 0.78, 4.45, 9.2, 1.4, colors, font)


def render_section(slide, prs: Presentation, data: dict[str, Any], colors: dict, font: str) -> None:
    add_accent_bar(slide, prs, colors)
    add_text(slide, data.get("kicker") or "Section", 0.75, 0.85, 2.5, 0.3, 13, colors["accent"], font, bold=True)
    add_text(slide, data.get("title") or "Section", 0.75, 1.55, 9.2, 1.0, 38, colors["primary"], font, bold=True)
    add_text(slide, data.get("narrative") or "", 0.8, 2.85, 7.2, 0.85, 20, colors["muted"], font)
    add_cards(slide, data.get("bullets") or [], 0.8, 4.15, 10.6, 1.6, colors, font)


def render_insight(
    slide,
    prs: Presentation,
    data: dict[str, Any],
    colors: dict,
    font: str,
    image_path: Path | None,
) -> None:
    add_header(slide, data, colors, font)
    if image_path:
        add_image_panel(slide, image_path, 7.35, 1.28, 4.95, 4.1, colors)
        text_width = 5.9
    else:
        text_width = 10.6
        add_visual_cards(slide, data.get("bullets") or [], 7.2, 3.9, 4.8, 1.35, colors, font)
    add_text(slide, data.get("narrative") or "", 0.75, 1.75, text_width, 0.8, 22, colors["ink"], font, bold=True)
    add_bullets(slide, data.get("bullets") or [], 0.95, 2.85, text_width, 2.2, colors, font)
    evidence = str(data.get("evidence") or "").strip()
    if evidence:
        add_evidence(slide, evidence, 0.85, 5.65, 10.8, 0.65, colors, font)


def render_split(
    slide,
    prs: Presentation,
    data: dict[str, Any],
    colors: dict,
    font: str,
    image_path: Path | None,
) -> None:
    add_header(slide, data, colors, font)
    bullets = data.get("bullets") or []
    midpoint = max(1, len(bullets) // 2)
    add_cards(slide, bullets[:midpoint], 0.75, 2.0, 5.45, 3.2, colors, font, title="Now")
    if image_path:
        add_image_panel(slide, image_path, 6.75, 2.0, 5.45, 3.2, colors)
    else:
        add_cards(slide, bullets[midpoint:], 6.75, 2.0, 5.45, 3.2, colors, font, title="Next")
    evidence = str(data.get("evidence") or "").strip()
    if evidence:
        add_evidence(slide, evidence, 0.85, 5.85, 10.8, 0.55, colors, font)


def render_data_slide(slide, prs: Presentation, data: dict[str, Any], colors: dict, font: str) -> None:
    add_header(slide, data, colors, font)
    chart = data.get("chart") or {}
    labels = [str(item) for item in chart.get("labels") or []]
    values = [float(item) for item in chart.get("values") or [] if isinstance(item, (int, float))]
    if not labels or not values:
        render_insight(slide, prs, data, colors, font, None)
        return
    add_bar_chart(slide, labels[:6], values[:6], 0.9, 2.05, 6.7, 3.25, colors, font)
    add_text(slide, chart.get("caption") or data.get("narrative") or "", 8.1, 2.1, 3.85, 1.2, 22, colors["ink"], font, bold=True)
    add_bullets(slide, data.get("bullets") or [], 8.25, 3.55, 3.65, 1.8, colors, font)


def add_header(slide, data: dict[str, Any], colors: dict, font: str) -> None:
    add_text(slide, data.get("kicker") or "Insight", 0.75, 0.42, 3.0, 0.28, 11, colors["accent"], font, bold=True)
    add_text(slide, data.get("title") or "Untitled", 0.72, 0.78, 10.8, 0.72, 28, colors["primary"], font, bold=True)


def add_text(slide, text: Any, x: float, y: float, w: float, h: float, size: int, color: RGBColor, font: str, bold: bool = False) -> None:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = fit_text(str(text or ""), size)
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, colors: dict, font: str) -> None:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets[:5]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = fit_text(str(bullet), 17)
        paragraph.level = 0
        paragraph.font.name = font
        paragraph.font.size = Pt(17)
        paragraph.font.color.rgb = colors["ink"]
        paragraph.space_after = Pt(8)


def add_cards(slide, items: list[str], x: float, y: float, w: float, h: float, colors: dict, font: str, title: str | None = None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    if title:
        add_text(slide, title, x, y - 0.42, w, 0.3, 14, colors["accent"], font, bold=True)
    count = max(1, min(3, len(items)))
    card_w = (w - 0.18 * (count - 1)) / count
    for index, item in enumerate(items[:count]):
        left = x + index * (card_w + 0.18)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(y), Inches(card_w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors["white"]
        shape.line.color.rgb = RGBColor(226, 221, 211)
        add_text(slide, item, left + 0.18, y + 0.18, card_w - 0.36, h - 0.25, 14, colors["ink"], font, bold=True)


def add_visual_cards(slide, items: list[str], x: float, y: float, w: float, h: float, colors: dict, font: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    for index, item in enumerate(items[:3]):
        top = y + index * (h + 0.18)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(226, 221, 211)
        add_text(slide, f"{index + 1}", x + 0.2, top + 0.22, 0.45, 0.35, 15, colors["accent"], font, bold=True)
        add_text(slide, item, x + 0.75, top + 0.18, w - 1.0, h - 0.25, 13, colors["ink"], font)


def add_evidence(slide, evidence: str, x: float, y: float, w: float, h: float, colors: dict, font: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 251, 240)
    shape.line.color.rgb = RGBColor(231, 217, 184)
    add_text(slide, f"Evidence: {evidence}", x + 0.2, y + 0.14, w - 0.4, h - 0.12, 12, colors["muted"], font)


def add_image_panel(slide, path: Path, x: float, y: float, w: float, h: float, colors: dict) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(226, 221, 211)
    slide.shapes.add_picture(str(path), Inches(x + 0.12), Inches(y + 0.12), width=Inches(w - 0.24), height=Inches(h - 0.24))


def add_bar_chart(slide, labels: list[str], values: list[float], x: float, y: float, w: float, h: float, colors: dict, font: str) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    max_value = max(values) if values else 1
    bar_gap = 0.13
    bar_w = (w - bar_gap * (len(values) - 1)) / max(1, len(values))
    for index, value in enumerate(values):
        height = max(0.15, h * (value / max_value))
        left = x + index * (bar_w + bar_gap)
        top = y + h - height
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(bar_w), Inches(height))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors["accent"] if index == len(values) - 1 else colors["primary"]
        bar.line.fill.background()
        add_text(slide, labels[index], left, y + h + 0.12, bar_w, 0.35, 9, colors["muted"], font)
        add_text(slide, f"{value:g}", left, top - 0.34, bar_w, 0.25, 10, colors["ink"], font, bold=True)


def add_accent_bar(slide, prs: Presentation, colors: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["accent"]
    shape.line.fill.background()


def add_big_number(slide, prs: Presentation, text: str, colors: dict, font: str) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(0.95), Inches(3.0), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["primary"]
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(9.2), Inches(1.95), Inches(3.0), Inches(0.6))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = font
    paragraph.font.size = Pt(30)
    paragraph.font.bold = True
    paragraph.font.color.rgb = colors["white"]


def add_page_number(slide, prs: Presentation, index: int, total: int, colors: dict, font: str) -> None:
    add_text(slide, f"{index:02d}/{total:02d}", 11.55, 7.02, 1.0, 0.22, 9, colors["muted"], font)


def has_chart(data: dict[str, Any]) -> bool:
    chart = data.get("chart") or {}
    return chart.get("type") in {"bar", "line"} and chart.get("labels") and chart.get("values")


def parse_hex(value: Any, fallback: RGBColor) -> RGBColor:
    from pptx.dml.color import RGBColor

    if not isinstance(value, str):
        return fallback
    match = re.fullmatch(r"#?([0-9a-fA-F]{6})", value.strip())
    if not match:
        return fallback
    raw = match.group(1)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def fit_text(value: str, target_size: int) -> str:
    limit = 130 if target_size >= 20 else 190
    value = re.sub(r"\s+", " ", value).strip()
    if len(value) <= limit:
        return value
    return textwrap.shorten(value, width=limit, placeholder="...")


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n[poc] cancelled", file=sys.stderr)
        raise SystemExit(130)
