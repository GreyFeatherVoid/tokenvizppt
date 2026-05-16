import re
from dataclasses import dataclass
from pathlib import Path
from urllib.parse import urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from app.core.settings import get_settings
from app.services.asset_store import AssetNotFoundError, get_asset_store
from app.services.deck_spec import prepare_elements_for_rendering
from app.services.db_session_repository import get_db_session_repository
from app.services.session_store import SessionNotFoundError, get_session_store, safe_session_id


@dataclass
class ExportedPptx:
    path: Path
    url: str


async def export_session_from_spec(session_id: str) -> ExportedPptx:
    session_id = safe_session_id(session_id)
    session = _load_session(session_id)
    slides = session.get("slides") or []
    if not slides or not all(slide.get("spec") for slide in slides):
        raise SessionNotFoundError("Slide specs are not available")

    export_root = get_settings().storage_root / "exports" / session_id
    export_root.mkdir(parents=True, exist_ok=True)
    output_path = export_root / f"{session_id}.pptx"

    prs = Presentation()
    first_spec = slides[0]["spec"]
    width = float(first_spec.get("size", {}).get("width") or 13.333)
    height = float(first_spec.get("size", {}).get("height") or 7.5)
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    prs.core_properties.title = str(session.get("topic") or "tokenvizPPT")
    prs.core_properties.author = "tokenvizPPT"
    blank = prs.slide_layouts[6]

    for slide_payload in slides:
        spec = slide_payload["spec"]
        ppt_slide = prs.slides.add_slide(blank)
        _set_background(ppt_slide, str(spec.get("background") or "#FFFFFF"))
        for element in prepare_elements_for_rendering(spec.get("elements") or []):
            if element.get("kind") == "shape":
                _add_shape(ppt_slide, element)
            elif element.get("kind") == "text":
                _add_text(ppt_slide, element)
            elif element.get("kind") == "image":
                _add_image(ppt_slide, element)

    prs.save(output_path)
    return ExportedPptx(
        path=output_path,
        url=f"/api/exports/{session_id}/{output_path.name}",
    )


def _load_session(session_id: str) -> dict:
    try:
        return get_db_session_repository().get_session_detail(session_id)
    except SessionNotFoundError:
        return get_session_store().get_session(session_id)


def _set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color, "FFFFFF")


def _add_shape(slide, element: dict) -> None:
    box = _box(element)
    shape_type = (
        MSO_SHAPE.ROUNDED_RECTANGLE
        if float(element.get("radius") or 0) > 0
        else MSO_SHAPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(box["x"]),
        Inches(box["y"]),
        Inches(box["w"]),
        Inches(box["h"]),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(str(element.get("fill") or "#FFFFFF"), "FFFFFF")
    border = str(element.get("border") or "")
    if border:
        shape.line.color.rgb = _rgb(border, "000000")
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()


def _add_text(slide, element: dict) -> None:
    text = str(element.get("text") or "")
    if not text.strip():
        return
    box = _box(element)
    textbox = slide.shapes.add_textbox(
        Inches(box["x"]),
        Inches(box["y"]),
        Inches(box["w"]),
        Inches(box["h"]),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = float(element.get("lineHeight") or 1.16)
    run = paragraph.add_run()
    run.text = _wrap_text(
        text,
        width_in=box["w"],
        font_size_pt=float(element.get("fontSize") or 12),
    )
    font = run.font
    font.size = Pt(float(element.get("fontSize") or 12))
    font.name = str(element.get("fontFace") or "Aptos")
    font.bold = bool(element.get("bold"))
    font.color.rgb = _rgb(str(element.get("color") or "#111827"), "111827")


def _add_image(slide, element: dict) -> None:
    path = _resolve_asset_path(str(element.get("src") or ""))
    if not path:
        return
    box = _box(element)
    slide.shapes.add_picture(
        str(path),
        Inches(box["x"]),
        Inches(box["y"]),
        width=Inches(box["w"]),
        height=Inches(box["h"]),
    )


def _box(element: dict) -> dict[str, float]:
    return {
        "x": max(0, float(element.get("x") or 0)),
        "y": max(0, float(element.get("y") or 0)),
        "w": max(0.01, float(element.get("w") or 0)),
        "h": max(0.01, float(element.get("h") or 0)),
    }


def _wrap_text(text: str, width_in: float, font_size_pt: float) -> str:
    if "\n" in text:
        return text
    has_cjk = bool(re.search(r"[\u3400-\u9fff\uf900-\ufaff]", text))
    capacity = _line_capacity(width_in, font_size_pt, has_cjk=has_cjk)
    if _visual_units(text, has_cjk=has_cjk) <= capacity:
        return text
    if has_cjk:
        return _wrap_by_visual_units(text, capacity, has_cjk=True)
    return _wrap_latin_text(text, capacity)


def _line_capacity(width_in: float, font_size_pt: float, *, has_cjk: bool) -> float:
    chars_per_inch_at_12pt = 5.15 if has_cjk else 9.2
    return max(3.0, width_in * chars_per_inch_at_12pt * 12 / max(font_size_pt, 1))


def _wrap_by_visual_units(text: str, capacity: float, *, has_cjk: bool) -> str:
    lines: list[str] = []
    current = ""
    current_units = 0.0
    for char in text:
        units = _char_units(char, has_cjk=has_cjk)
        if current and current_units + units > capacity:
            lines.append(current.rstrip())
            current = char.lstrip()
            current_units = _visual_units(current, has_cjk=has_cjk)
        else:
            current += char
            current_units += units
    if current:
        lines.append(current.rstrip())
    return "\n".join(lines)


def _wrap_latin_text(text: str, capacity: float) -> str:
    lines: list[str] = []
    current = ""
    current_units = 0.0
    for word in text.split(" "):
        token = word if not current else f" {word}"
        units = _visual_units(token, has_cjk=False)
        if current and current_units + units > capacity:
            lines.append(current)
            current = word
            current_units = _visual_units(word, has_cjk=False)
        else:
            current += token
            current_units += units
    if current:
        lines.append(current)
    return "\n".join(lines)


def _visual_units(text: str, *, has_cjk: bool) -> float:
    return sum(_char_units(char, has_cjk=has_cjk) for char in text)


def _char_units(char: str, *, has_cjk: bool) -> float:
    if char.isspace():
        return 0.35
    if re.match(r"[\u3400-\u9fff\uf900-\ufaff]", char):
        return 1.0
    if char in "，。；：！？、“”‘’（）《》—…":
        return 0.75
    return 0.55 if has_cjk else 1.0


def _resolve_asset_path(src: str) -> Path | None:
    path = urlparse(src).path
    match = re.search(r"/api/assets/([^/]+)/file$", path)
    if not match:
        return None
    try:
        return get_asset_store().get_asset_file_path(match.group(1))
    except (AssetNotFoundError, SessionNotFoundError):
        return None


def _rgb(value: str, fallback: str) -> RGBColor:
    raw = value.strip().replace("#", "").upper()
    if re.fullmatch(r"[0-9A-F]{3}", raw):
        raw = "".join(char * 2 for char in raw)
    if not re.fullmatch(r"[0-9A-F]{6}", raw):
        raw = fallback
    return RGBColor(
        int(raw[0:2], 16),
        int(raw[2:4], 16),
        int(raw[4:6], 16),
    )
