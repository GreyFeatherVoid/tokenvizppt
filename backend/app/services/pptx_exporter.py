# ruff: noqa: E501
import re
from dataclasses import dataclass
from pathlib import Path
from urllib.parse import urlparse

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from app.core.settings import get_settings
from app.services.asset_store import AssetNotFoundError, get_asset_store
from app.services.db_session_repository import get_db_session_repository
from app.services.session_store import SessionNotFoundError, get_session_store, safe_session_id

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
VIEWPORT_WIDTH = 1600
VIEWPORT_HEIGHT = 900


@dataclass
class ExportedPptx:
    path: Path
    url: str


async def export_session_to_editable_pptx(session_id: str) -> ExportedPptx:
    session_id = safe_session_id(session_id)
    session = _load_session(session_id)
    export_root = get_settings().storage_root / "exports" / session_id
    export_root.mkdir(parents=True, exist_ok=True)
    output_path = export_root / f"{session_id}.pptx"

    slides = []
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        try:
            page = await browser.new_page(
                viewport={"width": VIEWPORT_WIDTH, "height": VIEWPORT_HEIGHT},
                device_scale_factor=1,
            )
            for slide in session["slides"]:
                html = _prepare_html(slide["html"])
                await page.set_content(html, wait_until="networkidle")
                await page.wait_for_timeout(200)
                slides.append(await page.evaluate(_extract_script()))
        finally:
            await browser.close()

    _write_pptx(output_path, session["topic"], slides)
    return ExportedPptx(
        path=output_path,
        url=f"/api/exports/{session_id}/{output_path.name}",
    )


def get_export_file(session_id: str, file_name: str) -> Path:
    session_id = safe_session_id(session_id)
    if file_name != f"{session_id}.pptx":
        raise SessionNotFoundError("Export file not found")
    path = (get_settings().storage_root / "exports" / session_id / file_name).resolve()
    allowed_root = (get_settings().storage_root / "exports" / session_id).resolve()
    if allowed_root not in [path, *path.parents] or not path.exists():
        raise SessionNotFoundError("Export file not found")
    return path


def _load_session(session_id: str) -> dict:
    try:
        return get_db_session_repository().get_session_detail(session_id)
    except SessionNotFoundError:
        return get_session_store().get_session(session_id)


def _prepare_html(html: str) -> str:
    base = get_settings().public_base_url.rstrip("/")
    if "<base " in html.lower():
        return html
    return re.sub(r"(<head[^>]*>)", rf'\1<base href="{base}">', html, count=1, flags=re.I)


def _extract_script() -> str:
    return f"""
() => {{
  const slideWidthIn = {SLIDE_WIDTH_IN};
  const slideHeightIn = {SLIDE_HEIGHT_IN};
  const hexToRgb = (hex) => {{
    const raw = String(hex || '').replace('#', '').trim();
    if (!/^[0-9A-F]{{6}}$/i.test(raw)) return {{ r: 255, g: 255, b: 255 }};
    return {{
      r: Number.parseInt(raw.slice(0, 2), 16),
      g: Number.parseInt(raw.slice(2, 4), 16),
      b: Number.parseInt(raw.slice(4, 6), 16),
    }};
  }};
  const rgbToHexValue = (rgb) => [rgb.r, rgb.g, rgb.b]
    .map((part) => Math.max(0, Math.min(255, Math.round(Number(part) || 0))).toString(16).padStart(2, '0'))
    .join('')
    .toUpperCase();
  const blend = (fg, bg) => rgbToHexValue({{
    r: fg.r * fg.a + bg.r * (1 - fg.a),
    g: fg.g * fg.a + bg.g * (1 - fg.a),
    b: fg.b * fg.a + bg.b * (1 - fg.a),
  }});
  const parseCssColor = (value) => {{
    const source = String(value || '').trim();
    if (!source || source === 'transparent') return null;
    if (source.startsWith('#')) {{
      const raw = source.slice(1).toUpperCase();
      const hex = raw.length === 3 ? raw.split('').map((part) => part + part).join('') : raw;
      if (!/^[0-9A-F]{{6}}$/.test(hex)) return null;
      const rgb = hexToRgb(hex);
      return {{ ...rgb, a: 1 }};
    }}
    const match = source.match(/rgba?\\(\\s*(\\d+(?:\\.\\d+)?)(?:\\s*,\\s*|\\s+)(\\d+(?:\\.\\d+)?)(?:\\s*,\\s*|\\s+)(\\d+(?:\\.\\d+)?)(?:\\s*(?:,|\\/)\\s*(\\d+(?:\\.\\d+)?%?))?/i);
    if (!match) return null;
    const alpha = match[4] === undefined ? 1 : String(match[4]).endsWith('%') ? Number.parseFloat(match[4]) / 100 : Number(match[4]);
    if (alpha <= 0.04) return null;
    return {{
      r: Math.max(0, Math.min(255, Number(match[1]) || 0)),
      g: Math.max(0, Math.min(255, Number(match[2]) || 0)),
      b: Math.max(0, Math.min(255, Number(match[3]) || 0)),
      a: Math.max(0, Math.min(1, alpha)),
    }};
  }};
  const rgbToHex = (value, fallback = 'FFFFFF') => {{
    const color = parseCssColor(value);
    if (!color) return '';
    return color.a >= 0.995 ? rgbToHexValue(color) : blend(color, hexToRgb(fallback));
  }};
  const firstPaintColor = (backgroundImage) => {{
    const source = String(backgroundImage || '');
    if (!source || source === 'none') return '';
    const match = source.match(/#[0-9A-Fa-f]{{3,8}}|rgba?\\([^)]*\\)/);
    return match ? match[0] : '';
  }};
  const paintToHex = (backgroundColor, backgroundImage, fallback = 'FFFFFF') => {{
    const gradientColor = firstPaintColor(backgroundImage);
    return rgbToHex(gradientColor || backgroundColor, fallback) || fallback;
  }};
  const rootCandidates = Array.from(document.querySelectorAll('[data-ppt-page], .ppt-page, .ppt-page-root, .slide, article, main, section'));
  const root = rootCandidates
    .map((element) => ({{ element, rect: element.getBoundingClientRect() }}))
    .filter((item) => item.rect.width > 320 && item.rect.height > 180)
    .sort((a, b) => (b.rect.width * b.rect.height) - (a.rect.width * a.rect.height))[0]?.element || document.body;
  const rootRect = root.getBoundingClientRect();
  const rootStyle = window.getComputedStyle(root);
  const bodyStyle = window.getComputedStyle(document.body);
  const bodyBg = paintToHex(bodyStyle.backgroundColor, bodyStyle.backgroundImage, 'FFFFFF');
  const rootBg = paintToHex(rootStyle.backgroundColor, rootStyle.backgroundImage, bodyBg);
  const toBox = (element) => {{
    const rect = element.getBoundingClientRect();
    return {{
      rect,
      x: (rect.left - rootRect.left) / rootRect.width * slideWidthIn,
      y: (rect.top - rootRect.top) / rootRect.height * slideHeightIn,
      w: rect.width / rootRect.width * slideWidthIn,
      h: rect.height / rootRect.height * slideHeightIn,
    }};
  }};
  const visible = (element, style, rect) => {{
    if (!style || style.display === 'none' || style.visibility === 'hidden') return false;
    if (Number(style.opacity || '1') < 0.04) return false;
    if (rect.width < 2 || rect.height < 2) return false;
    if (element.closest('script, style, noscript, iframe, object, embed')) return false;
    return true;
  }};
  const texts = [];
  const consumedTextElements = [];
  const textNodes = Array.from(root.querySelectorAll('h1,h2,h3,h4,h5,h6,p,li,blockquote,td,th,figcaption,div[data-edit-id],span[data-edit-id],small[data-edit-id],strong[data-edit-id],em[data-edit-id]'));
  const hasTextChild = (element) => Boolean(element.querySelector('h1,h2,h3,h4,h5,h6,p,li,blockquote,td,th,figcaption,div[data-edit-id],span[data-edit-id],small[data-edit-id],strong[data-edit-id],em[data-edit-id]'));
  const insideConsumed = (element) => consumedTextElements.some((parent) => parent !== element && parent.contains(element));
  const lineTextsForElement = (element) => {{
    const walker = document.createTreeWalker(element, NodeFilter.SHOW_TEXT, {{
      acceptNode: (node) => node.nodeValue && node.nodeValue.trim()
        ? NodeFilter.FILTER_ACCEPT
        : NodeFilter.FILTER_REJECT,
    }});
    const lines = [];
    while (walker.nextNode()) {{
      const node = walker.currentNode;
      const raw = node.nodeValue || '';
      for (let start = 0; start < raw.length; start += 1) {{
        if (/\\s/.test(raw[start])) continue;
        const range = document.createRange();
        range.setStart(node, start);
        range.setEnd(node, start + 1);
        const firstRect = range.getClientRects()[0];
        range.detach();
        if (!firstRect) continue;
        let end = start + 1;
        for (; end < raw.length; end += 1) {{
          const probe = document.createRange();
          probe.setStart(node, end);
          probe.setEnd(node, end + 1);
          const rect = probe.getClientRects()[0];
          probe.detach();
          if (!rect || Math.abs(rect.top - firstRect.top) > Math.max(2, firstRect.height * 0.45)) break;
        }}
        const text = raw.slice(start, end).replace(/\\s+/g, ' ').trim();
        if (text) {{
          const lineRange = document.createRange();
          lineRange.setStart(node, start);
          lineRange.setEnd(node, end);
          const rects = Array.from(lineRange.getClientRects())
            .filter((rect) => rect.width > 1 && rect.height > 1);
          lineRange.detach();
          if (rects.length) {{
            const left = Math.min(...rects.map((rect) => rect.left));
            const top = Math.min(...rects.map((rect) => rect.top));
            const right = Math.max(...rects.map((rect) => rect.right));
            const bottom = Math.max(...rects.map((rect) => rect.bottom));
            lines.push({{
              text,
              rect: {{ left, top, right, bottom, width: right - left, height: bottom - top }},
            }});
          }}
        }}
        start = Math.max(start, end - 1);
      }}
    }}
    return lines;
  }};
  const seen = new Set();
  for (const element of textNodes) {{
    if (texts.length >= 120) break;
    if (insideConsumed(element)) continue;
    const style = window.getComputedStyle(element);
    const box = toBox(element);
    if (!visible(element, style, box.rect)) continue;
    if (hasTextChild(element) && !element.dataset.editId) continue;
    const visualLines = lineTextsForElement(element);
    if (!visualLines.length) continue;
    for (const line of visualLines) {{
      if (texts.length >= 160) break;
      if (!line.text || line.text.length > 1200) continue;
      const lineBox = {{
        x: (line.rect.left - rootRect.left) / rootRect.width * slideWidthIn,
        y: (line.rect.top - rootRect.top) / rootRect.height * slideHeightIn,
        w: line.rect.width / rootRect.width * slideWidthIn,
        h: line.rect.height / rootRect.height * slideHeightIn,
      }};
      const key = [line.text.replace(/\\s+/g, ' '), Math.round(line.rect.left), Math.round(line.rect.top), Math.round(line.rect.width), Math.round(line.rect.height)].join('|');
      if (seen.has(key)) continue;
      seen.add(key);
      texts.push({{
        text: line.text,
        x: lineBox.x,
        y: lineBox.y,
        w: Math.max(lineBox.w, box.w * 0.2),
        h: Math.max(lineBox.h, box.h * 0.08),
        fontSize: (Number.parseFloat(style.fontSize || '16') || 16) * 72 / (rootRect.height / slideHeightIn),
        fontFace: String(style.fontFamily || '').split(',')[0].replace(/[\\\"']/g, '').trim(),
        color: rgbToHex(style.color, rootBg) || '111827',
        bold: (Number.parseInt(style.fontWeight || '400', 10) || 400) >= 600 || /^H[1-6]$/i.test(element.tagName),
        italic: style.fontStyle === 'italic' || style.fontStyle === 'oblique',
        align: ['center', 'right', 'justify'].includes(style.textAlign) ? style.textAlign : 'left',
        lineHeight: Number.parseFloat(style.lineHeight || '') || 0,
      }});
    }}
    consumedTextElements.push(element);
  }}
  const shapes = [];
  const shapeNodes = Array.from(root.querySelectorAll('section,main,article,header,footer,aside,div,figure,table,td,th'));
  for (const element of shapeNodes) {{
    if (shapes.length >= 80 || element === root) break;
    const style = window.getComputedStyle(element);
    const box = toBox(element);
    if (!visible(element, style, box.rect)) continue;
    const fill = paintToHex(style.backgroundColor, style.backgroundImage, rootBg);
    const border = rgbToHex(style.borderColor, rootBg);
    const borderWidth = Number.parseFloat(style.borderWidth || '0') || 0;
    if ((!fill || fill === rootBg) && (!border || borderWidth <= 0 || style.borderStyle === 'none')) continue;
    if (box.rect.width * box.rect.height < rootRect.width * rootRect.height * 0.003) continue;
    shapes.push({{
      x: box.x,
      y: box.y,
      w: box.w,
      h: box.h,
      fill,
      border,
      borderWidth,
      radius: Number.parseFloat(style.borderRadius || style.borderTopLeftRadius || '0') || 0,
    }});
  }}
  const images = [];
  for (const element of Array.from(root.querySelectorAll('img'))) {{
    if (images.length >= 40) break;
    const style = window.getComputedStyle(element);
    const box = toBox(element);
    if (!visible(element, style, box.rect)) continue;
    images.push({{
      src: element.currentSrc || element.src || '',
      alt: element.getAttribute('alt') || '',
      x: box.x,
      y: box.y,
      w: box.w,
      h: box.h,
    }});
  }}
  return {{
    backgroundColor: rootBg || bodyBg || 'FFFFFF',
    texts,
    shapes,
    images,
  }};
}}
"""


def _write_pptx(output_path: Path, title: str, slides: list[dict]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]
    prs.core_properties.title = title
    prs.core_properties.author = "tokenvizPPT"

    for source in slides:
        slide = prs.slides.add_slide(blank)
        _set_background(slide, source.get("backgroundColor") or "FFFFFF")
        for shape in source.get("shapes") or []:
            _add_shape(slide, shape)
        for image in source.get("images") or []:
            _add_image(slide, image)
        for text in source.get("texts") or []:
            _add_text(slide, text)

    prs.save(output_path)


def _set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color, "FFFFFF")


def _add_shape(slide, shape: dict) -> None:
    box = _box(shape)
    if box["w"] <= 0 or box["h"] <= 0:
        return
    ppt_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if float(shape.get("radius") or 0) > 1 else MSO_SHAPE.RECTANGLE,
        Inches(box["x"]),
        Inches(box["y"]),
        Inches(box["w"]),
        Inches(box["h"]),
    )
    fill = str(shape.get("fill") or "")
    if fill:
        ppt_shape.fill.solid()
        ppt_shape.fill.fore_color.rgb = _rgb(fill, "FFFFFF")
    else:
        ppt_shape.fill.background()
    border = str(shape.get("border") or "")
    if border and float(shape.get("borderWidth") or 0) > 0:
        ppt_shape.line.color.rgb = _rgb(border, "000000")
        ppt_shape.line.width = Pt(float(shape.get("borderWidth") or 1) * 0.75)
    else:
        ppt_shape.line.fill.background()


def _add_text(slide, text: dict) -> None:
    value = str(text.get("text") or "").strip()
    if not value:
        return
    box = _box(text)
    if box["w"] <= 0 or box["h"] <= 0:
        return
    font_size = _safe_font_size(value, box["w"], box["h"], float(text.get("fontSize") or 14))
    safe_height = _safe_text_height(value, box["w"], box["h"], font_size)
    textbox = slide.shapes.add_textbox(
        Inches(box["x"]),
        Inches(box["y"]),
        Inches(box["w"]),
        Inches(safe_height),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    lines = _text_lines(value)
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = _alignment(str(text.get("align") or "left"))
        paragraph.line_spacing = 1.08
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = str(text.get("fontFace") or "Aptos")
        run.font.color.rgb = _rgb(str(text.get("color") or "111827"), "111827")
        run.font.bold = bool(text.get("bold"))
        run.font.italic = bool(text.get("italic"))


def _add_image(slide, image: dict) -> None:
    path = _resolve_image_path(str(image.get("src") or ""))
    if not path or not path.exists():
        return
    box = _box(image)
    if box["w"] <= 0 or box["h"] <= 0:
        return
    slide.shapes.add_picture(
        str(path),
        Inches(box["x"]),
        Inches(box["y"]),
        width=Inches(box["w"]),
        height=Inches(box["h"]),
    )


def _resolve_image_path(src: str) -> Path | None:
    path = urlparse(src).path
    match = re.search(r"/api/assets/([^/]+)/file$", path)
    if not match:
        return None
    try:
        return get_asset_store().get_asset_file_path(match.group(1))
    except AssetNotFoundError:
        return None


def _box(item: dict) -> dict[str, float]:
    return {
        "x": _clamp(float(item.get("x") or 0), 0, SLIDE_WIDTH_IN),
        "y": _clamp(float(item.get("y") or 0), 0, SLIDE_HEIGHT_IN),
        "w": _clamp(float(item.get("w") or 0), 0, SLIDE_WIDTH_IN),
        "h": _clamp(float(item.get("h") or 0), 0, SLIDE_HEIGHT_IN),
    }


def _rgb(value: str, fallback: str) -> RGBColor:
    raw = value.strip().replace("#", "").upper()
    if re.fullmatch(r"[0-9A-F]{3}", raw):
        raw = "".join(char * 2 for char in raw)
    if not re.fullmatch(r"[0-9A-F]{6}", raw):
        raw = fallback
    return RGBColor(
        int(raw[0:2], 16),
        int(raw[2:4], 16),
        int(raw[4:6], 16),
    )


def _alignment(value: str):
    from pptx.enum.text import PP_ALIGN

    if value == "center":
        return PP_ALIGN.CENTER
    if value == "right":
        return PP_ALIGN.RIGHT
    if value == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def _clamp(value: float, low: float, high: float) -> float:
    return max(low, min(high, value))


def _safe_text_height(text: str, width_in: float, dom_height_in: float, font_size_pt: float) -> float:
    if width_in <= 0:
        return max(dom_height_in, 0.18)
    has_cjk = bool(re.search(r"[\u3400-\u9fff\uf900-\ufaff]", text))
    chars_per_line = max(4, int(width_in * (8.2 if has_cjk else 11.5) * 12 / max(font_size_pt, 1)))
    explicit_lines = _text_lines(text)
    line_count = sum(max(1, (len(line) + chars_per_line - 1) // chars_per_line) for line in explicit_lines)
    estimated = line_count * font_size_pt / 72 * 1.24 + 0.06
    return _clamp(max(dom_height_in * 1.18, estimated, 0.16), 0.12, SLIDE_HEIGHT_IN)


def _safe_font_size(text: str, width_in: float, height_in: float, font_size_pt: float) -> float:
    size = _clamp(font_size_pt, 6, 96)
    for _ in range(14):
        if _safe_text_height(text, width_in, height_in, size) <= max(height_in * 1.5, height_in + 0.12):
            return size
        next_size = size - 1
        if next_size < max(6, font_size_pt * 0.72):
            return max(6, next_size)
        size = next_size
    return size


def _text_lines(value: str) -> list[str]:
    lines = [
        re.sub(r"[^\S\n]+", " ", line).strip()
        for line in value.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    ]
    return [line for line in lines if line] or [value.strip()]
